from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation

from native_pack_pdf_export import (
    export_selected_pack_questions_to_pdf,
    load_native_slide_groups,
    match_selected_slides,
    parse_selected_slides,
    text_similarity,
)


FIXTURE_ROOT = Path(__file__).resolve().parents[1] / "contracts" / "quail-ultra-qbank" / "v1" / "fixtures"


def _write_pptx(path: Path, slide_texts: list[str]) -> None:
    presentation = Presentation()
    blank = presentation.slide_layouts[6]
    for text in slide_texts:
        slide = presentation.slides.add_slide(blank)
        textbox = slide.shapes.add_textbox(914400, 914400, 7315200, 1828800)
        textbox.text_frame.text = text
    presentation.save(path)


def test_text_similarity_uses_overlap_and_sequence() -> None:
    left = "A neonate has jaundice at 18 hours of life. Which diagnosis?"
    right = "A neonate has jaundice at 18 hours of life and needs evaluation."
    unrelated = "A toddler has recurrent otitis media and eczema."

    assert text_similarity(left, right) > 0.55
    assert text_similarity(left, unrelated) < text_similarity(left, right)


def test_load_native_slide_groups_groups_questions_by_source_slide() -> None:
    groups = load_native_slide_groups(FIXTURE_ROOT / "native-pack-minimal")

    assert sorted(groups) == ["peds-sample:1", "peds-sample:2", "peds-sample:3"]
    assert groups["peds-sample:1"].questions[0].qid == "peds.sample.s001.q01"
    assert groups["peds-sample:1"].source_slide_path is not None


def test_match_selected_slides_finds_existing_native_questions(tmp_path) -> None:
    pptx = tmp_path / "selected.pptx"
    _write_pptx(
        pptx,
        [
            "A neonate has jaundice at 18 hours of life. Phototherapy and neonatal jaundice.",
            "A toddler has recurrent otitis media and eczema. Primary immunodeficiency.",
        ],
    )
    groups = load_native_slide_groups(FIXTURE_ROOT / "native-pack-minimal")
    selected = parse_selected_slides(pptx, tmp_path / "work")

    matches = match_selected_slides(selected, groups, min_score=0.20)

    assert [match.matched_source_key for match in matches] == ["peds-sample:2", "peds-sample:3"]
    assert matches[0].question_ids == ["peds.sample.s002.q01"]
    assert all(match.confidence in {"matched", "ambiguous"} for match in matches)


def test_export_selected_pack_questions_to_pdf_writes_reports_and_pdf(tmp_path) -> None:
    pptx = tmp_path / "selected.pptx"
    output_pdf = tmp_path / "matched.pdf"
    _write_pptx(
        pptx,
        [
            "A 6-year-old has fever, cough, tachypnea, and focal crackles.",
            "A neonate has jaundice at 18 hours of life.",
        ],
    )

    matches, json_report, csv_report = export_selected_pack_questions_to_pdf(
        pack_dir=FIXTURE_ROOT / "native-pack-minimal",
        pptx_path=pptx,
        output_pdf=output_pdf,
        output_dir=tmp_path / "reports",
        min_score=0.20,
    )

    assert output_pdf.exists()
    assert output_pdf.stat().st_size > 0
    assert json_report.exists()
    assert csv_report.exists()
    report = json.loads(json_report.read_text(encoding="utf-8"))
    assert [item["matchedSourceKey"] for item in report] == ["peds-sample:1", "peds-sample:2"]
    assert [match.question_ids for match in matches] == [["peds.sample.s001.q01"], ["peds.sample.s002.q01"]]
