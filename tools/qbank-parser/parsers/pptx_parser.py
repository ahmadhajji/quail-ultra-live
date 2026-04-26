"""
PPTX Parser Module

Extracts content from PowerPoint (.pptx) files:
- All text from each slide
- Speaker notes
- Highlighted/colored text (potential correct answers)
- Embedded images
- Slide numbers for tracking
"""

from pathlib import Path
from typing import Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
import io
import os

from domain.models import SlideContent
from utils.image_filters import is_placeholder_image_for_extraction


def is_yellow_like(color: RGBColor) -> bool:
    """Check if a color is yellow or yellow-like (highlighting)."""
    if color is None:
        return False
    
    r, g, b = color[0], color[1], color[2]
    
    # Yellow-like colors have high R and G, low B
    # We're lenient here to catch various yellow shades
    if r > 200 and g > 180 and b < 150:
        return True
    
    # Also check for common highlight yellows
    # Standard yellow: RGB(255, 255, 0)
    # Light yellow: RGB(255, 255, 153)
    # Google yellow: RGB(255, 242, 0) or RGB(255, 235, 59)
    if r > 230 and g > 200 and b < 180:
        return True
    
    return False


def extract_text_from_shape(shape, slide_content: SlideContent) -> None:
    """Extract text from a shape, detecting highlighting."""
    if not shape.has_text_frame:
        return
    
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            text = run.text.strip()
            if not text:
                continue
            
            slide_content.texts.append(text)
            
            # Check for highlighting
            is_highlighted = False
            
            # Check fill color of the shape
            try:
                if hasattr(shape, 'fill') and shape.fill.type is not None:
                    if shape.fill.fore_color and shape.fill.fore_color.type == 1:  # RGB
                        if is_yellow_like(shape.fill.fore_color.rgb):
                            is_highlighted = True
            except Exception:
                pass
            
            # Check text highlight/background color
            try:
                font = run.font
                if hasattr(font, 'highlight_color') and font.highlight_color is not None:
                    is_highlighted = True
            except Exception:
                pass
            
            # Check font color (sometimes correct answers are in specific colors)
            try:
                font = run.font
                if font.color and font.color.type == 1:  # RGB
                    rgb = font.color.rgb
                    if is_yellow_like(rgb):
                        is_highlighted = True
            except Exception:
                pass
            
            if is_highlighted:
                slide_content.highlighted_texts.append(text)


def extract_images_from_slide(slide, slide_number: int, output_dir: Path) -> list[str]:
    """Extract embedded images from a slide."""
    images = []
    img_counter = 0
    
    for shape in slide.shapes:
        if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
            try:
                image = shape.image
                image_bytes = image.blob
                image_ext = image.ext
                
                img_filename = f"slide_{slide_number}_img_{img_counter}.{image_ext}"
                img_path = output_dir / img_filename
                
                with open(img_path, 'wb') as f:
                    f.write(image_bytes)

                if _is_placeholder_image(img_path):
                    # Skip known empty frame/white placeholder images early to reduce AI cost.
                    img_path.unlink(missing_ok=True)
                    continue

                images.append(str(img_path))
                img_counter += 1
            except Exception as e:
                print(f"Warning: Could not extract image from slide {slide_number}: {e}")
    
    return images


def _extract_slide_background_rgb(slide) -> tuple[int, int, int] | None:
    """Return slide background RGB when explicitly set."""
    try:
        fill = slide.background.fill
        fore_color = getattr(fill, "fore_color", None)
        rgb = getattr(fore_color, "rgb", None)
        if rgb is None:
            return None
        return int(rgb[0]), int(rgb[1]), int(rgb[2])
    except Exception:
        return None


def classify_slide_consensus(slide) -> str:
    """Best-effort mapping of slide background color to consensus state."""
    rgb = _extract_slide_background_rgb(slide)
    if rgb is None:
        return ""
    r, g, b = rgb
    if r >= 240 and g >= 240 and b >= 240:
        return "clear"
    if is_yellow_like(RGBColor(r, g, b)):
        return "no_consensus"
    if b >= max(r, g) + 20 and b >= 150:
        return "consensus"
    return ""


def _is_placeholder_image(image_path: Path) -> bool:
    """Return True if an image appears to be a white/empty placeholder."""
    return is_placeholder_image_for_extraction(image_path)


def extract_speaker_notes(slide) -> str:
    """Extract speaker notes from a slide."""
    try:
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            notes_text = notes_slide.notes_text_frame.text
            return notes_text.strip()
    except Exception:
        pass
    return ""


def parse_pptx(pptx_path: str | Path, output_dir: Optional[Path] = None) -> list[SlideContent]:
    """
    Parse a PowerPoint file and extract all content.
    
    Args:
        pptx_path: Path to the .pptx file
        output_dir: Directory to save extracted images (defaults to data/extracted/)
    
    Returns:
        List of SlideContent objects, one per slide
    """
    pptx_path = Path(pptx_path)
    
    if not pptx_path.exists():
        raise FileNotFoundError(f"PPTX file not found: {pptx_path}")
    
    if output_dir is None:
        output_dir = pptx_path.parent / "extracted"
    
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)
    
    prs = Presentation(pptx_path)
    slides_content = []
    
    for slide_number, slide in enumerate(prs.slides, start=1):
        slide_content = SlideContent(slide_number=slide_number)
        
        # Extract text from all shapes
        for shape in slide.shapes:
            extract_text_from_shape(shape, slide_content)
            
            # Store shape info for debugging
            try:
                shape_info = {
                    "type": str(shape.shape_type),
                    "name": shape.name if hasattr(shape, 'name') else "unknown",
                    "has_text": shape.has_text_frame if hasattr(shape, 'has_text_frame') else False,
                }
                slide_content.raw_shapes_info.append(shape_info)
            except Exception:
                pass
        
        # Extract speaker notes
        slide_content.speaker_notes = extract_speaker_notes(slide)
        
        # Extract images
        slide_content.images = extract_images_from_slide(slide, slide_number, output_dir)
        slide_content.slide_consensus_status = classify_slide_consensus(slide)

        slides_content.append(slide_content)
    
    return slides_content


def get_slide_summary(slide_content: SlideContent) -> str:
    """Get a brief summary of slide content for display."""
    summary_parts = []
    
    if slide_content.texts:
        # First non-empty text as title
        first_text = slide_content.texts[0][:80]
        if len(slide_content.texts[0]) > 80:
            first_text += "..."
        summary_parts.append(f"Text: {first_text}")
    
    if slide_content.highlighted_texts:
        summary_parts.append(f"Highlighted: {', '.join(slide_content.highlighted_texts[:3])}")
    
    if slide_content.speaker_notes:
        notes_preview = slide_content.speaker_notes[:50]
        if len(slide_content.speaker_notes) > 50:
            notes_preview += "..."
        summary_parts.append(f"Notes: {notes_preview}")
    
    if slide_content.images:
        summary_parts.append(f"Images: {len(slide_content.images)}")
    
    return " | ".join(summary_parts) if summary_parts else "Empty slide"


if __name__ == "__main__":
    # Quick test
    import sys
    if len(sys.argv) > 1:
        pptx_file = sys.argv[1]
        print(f"Parsing {pptx_file}...")
        
        slides = parse_pptx(pptx_file)
        
        for slide in slides[:5]:  # Show first 5
            print(f"\n--- Slide {slide.slide_number} ---")
            print(get_slide_summary(slide))
